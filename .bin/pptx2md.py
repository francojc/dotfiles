import os
import argparse
import logging
from pptx import Presentation
import shutil

# Set up logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)


def pptx_to_markdown(input_dir, output_dir):
    """Converts PPTX files to Markdown files."""

    logging.info(f"Processing files from input directory: {input_dir}")

    # Create the output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    # Keep track of the number of files processed
    file_num = 0

    for filename in os.listdir(input_dir):
        if filename.lower().endswith(".pptx"):
            pptx_file = os.path.join(input_dir, filename)
            output_md = os.path.join(output_dir, os.path.splitext(filename)[0] + ".md")
            image_dir = os.path.join(output_dir, "images")

            logging.info(f"Converting {pptx_file} to {output_md}")

            # Create the images folder if it doesn't exist
            os.makedirs(image_dir, exist_ok=True)

            file_num += 1

            try:
                # Load the presentation
                prs = Presentation(pptx_file)

                # Open the output markdown file
                with open(output_md, "w") as md_file:
                    # Loop through each slide in the presentation
                    for slide_num, slide in enumerate(prs.slides):
                        md_file.write(f"# Slide {slide_num + 1}\n\n")

                        # Extract and write all text from shapes that have text frames
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text = shape.text.strip()
                                if text:  # Avoid writing empty lines
                                    md_file.write(f"{text}\n\n")

                        # Extract images and save them
                        for shape in slide.shapes:
                            if hasattr(shape, "image") and shape.image:
                                # Extract image details
                                image = shape.image
                                image_bytes = image.blob
                                image_format = image.ext
                                # include the file number, slide number, and shape id in the image filename
                                image_filename = f"image_{file_num}_{slide_num + 1}_{shape.shape_id}.{image_format}"
                                image_path = os.path.join(image_dir, image_filename)

                                # Save the image to the images folder
                                with open(image_path, "wb") as img_file:
                                    img_file.write(image_bytes)

                                # Link the image in Markdown
                                md_file.write(
                                    f"![Image {slide_num + 1}]({os.path.join('images',image_filename)})\n\n"
                                )
                        md_file.write("---\n\n")  # Slide separator

                logging.info(
                    f"Conversion complete! Markdown saved to {output_md} and images are saved in the '{image_dir}' folder."
                )

            except Exception as e:
                logging.error(f"Error processing {pptx_file}: {e}")

    logging.info("Finished processing all files.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PPTX files to Markdown.")
    parser.add_argument(
        "-i",
        "--input",
        required=True,
        help="Path to the input directory containing PPTX files.",
    )
    parser.add_argument(
        "-o",
        "--output",
        required=True,
        help="Path to the output directory where Markdown files will be saved.",
    )
    args = parser.parse_args()

    pptx_to_markdown(args.input, args.output)
